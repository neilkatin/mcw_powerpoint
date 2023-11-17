#!/usr/bin/env python

import os
import re
import sys
import logging
import argparse
import datetime
import json
import io
import zipfile
import itertools
import functools

import O365
import lxml
import pptx

import config as config_static
import arc_o365
import init_config
import init_logging



NOW = datetime.datetime.now().astimezone()
DATESTAMP = NOW.strftime("%Y-%m-%d")
TIMESTAMP = NOW.strftime("%Y-%m-%d %H:%M:%S %Z")
FILESTAMP = NOW.strftime("%Y-%m-%d %H-%M-%S %Z")
EMAILSTAMP = NOW.strftime("%Y-%m-%d %H-%M")

POWERPOINT_SUFFIX = 'main powerpoint.pptx'
#DRIVE_PATH = "/Current Webinar/extra"
DRIVE_PATH = "/Current Webinar"


def main():
    args = parse_args()
    if args.debug:
        logging.getLogger().setLevel(logging.DEBUG)
    log.debug("running...")

    config = init_config.init_config(config_static, ".env")

    o365_obj = arc_o365.arc_o365(config, tenant_id=config.TENANT_ID)
    o365_account = o365_obj.get_account()
    if False:
        m = o365_account.new_message(resource='masscarewebinar@redcross.org')
        m.to.add('neil@askneil.com')
        m.subject = f"test_message"
        m.body = f"test body"
        m.send()

    storage = o365_account.storage()
    log.debug(f"storage: { storage }")

    drive = storage.get_drive(config.DRIVE_ID)
    log.debug(f"drive: { drive }")

    current_webinar = drive.get_item_by_path(DRIVE_PATH)
    log.debug(f"current_webinar: { current_webinar }")

    items = list(current_webinar.get_items())
    #log.debug(f"items: { items }")

    pptx_item = None
    for item in items:
        name = item.name.lower()
        #log.debug(f"examining '{ name }' against '{ POWERPOINT_SUFFIX }'")
        if name.endswith(POWERPOINT_SUFFIX):
            pptx_item = item
            break

    if pptx_item is not None:
        stream = io.BytesIO()
        result = pptx_item.download(output=stream)
        pptx_bytes = stream.getvalue()
        log.debug(f"downloaded item { item.name }.  result { result } len { len(stream.getvalue()) }")

        pptx_stream = io.BytesIO(initial_bytes=pptx_bytes)
        presentation = pptx.Presentation(pptx=pptx_stream)

        for slide_index, slide in enumerate(presentation.slides):
            if slide.has_notes_slide:
                notes = slide.notes_slide
                log.debug(f"slide { slide_index } text_frame.text { len(notes.notes_text_frame.text) } '{ notes.notes_text_frame.text }'")
                #element = notes.element
                #log.debug(f"         element tag { lxml.etree.tostring(element) }\n")
            else:
                notes = None
                log.debug(f"slide { slide_index } no notes")


    else:
        log.info(f"could not find powerpoint item")



def parse_args():
    parser = argparse.ArgumentParser(
            description="tools to support Disaster Transportation Tools reporting",
            allow_abbrev=False)
    parser.add_argument("--debug", help="turn on debugging output", action="store_true")

    args = parser.parse_args()
    
    return args



if __name__ == "__main__":
    init_logging.init_logging(__name__)
    log = logging.getLogger(__name__)
    main()
else:
    log = logging.getLogger(__name__)

